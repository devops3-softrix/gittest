"""PDF conversion and manipulation tools."""
import csv
import io
import os
import tempfile
import zipfile
from io import BytesIO

from django.http import HttpResponse
from django.shortcuts import render

from tools.utils import TempDir, file_response, zip_response
from tools.conversion_limit import check_conversion_limit, log_conversion


# ---------------------------------------------------------------------------
# Helper: check limit before conversion, log after
# ---------------------------------------------------------------------------

def _check_limit(request):
    """Return an error response if the user has exceeded their limit, else None."""
    allowed, remaining = check_conversion_limit(request)
    if not allowed:
        return render(request, "tools/error.html", {
            "error": "You have used all 3 free conversions. Please sign up for unlimited access."
        })
    return None


# ---------------------------------------------------------------------------
# PDF → DOCX
# ---------------------------------------------------------------------------

def pdf_to_docx(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            from pdf2docx import Converter
            uploaded = request.FILES["pdf_file"]
            with TempDir() as tmp:
                src = tmp / "input.pdf"
                src.write_bytes(uploaded.read())
                out = tmp / "output.docx"
                cv = Converter(str(src))
                cv.convert(str(out))
                cv.close()
                log_conversion(request, "pdf_to_docx")
                return file_response(
                    out.read_bytes(),
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "converted.docx",
                )
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → DOCX",
        "description": "Convert a PDF file into an editable Word document.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "📄",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → TXT
# ---------------------------------------------------------------------------

def pdf_to_txt(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            from pypdf import PdfReader
            uploaded = request.FILES["pdf_file"]
            reader = PdfReader(uploaded)
            text = "\n\n".join(
                page.extract_text() or "" for page in reader.pages
            )
            log_conversion(request, "pdf_to_txt")
            return file_response(text.encode("utf-8"), "text/plain", "extracted.txt")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → TXT",
        "description": "Extract all text content from a PDF file.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "📃",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → PNG  (each page → zipped PNGs)
# ---------------------------------------------------------------------------

def pdf_to_png(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            import fitz
            uploaded = request.FILES["pdf_file"]
            doc = fitz.open(stream=uploaded.read(), filetype="pdf")
            files = {}
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                files[f"page_{i+1}.png"] = pix.tobytes("png")
            doc.close()
            log_conversion(request, "pdf_to_png")
            return zip_response(files, "pages_png.zip")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → PNG",
        "description": "Convert each PDF page into a high-quality PNG image.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "🖼️",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → JPG
# ---------------------------------------------------------------------------

def pdf_to_jpg(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            import fitz
            uploaded = request.FILES["pdf_file"]
            doc = fitz.open(stream=uploaded.read(), filetype="pdf")
            files = {}
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                files[f"page_{i+1}.jpg"] = pix.tobytes("jpeg")
            doc.close()
            log_conversion(request, "pdf_to_jpg")
            return zip_response(files, "pages_jpg.zip")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → JPG",
        "description": "Convert each PDF page into a JPEG image.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "🖼️",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → BMP
# ---------------------------------------------------------------------------

def pdf_to_bmp(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            import fitz
            from PIL import Image
            uploaded = request.FILES["pdf_file"]
            doc = fitz.open(stream=uploaded.read(), filetype="pdf")
            files = {}
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                buf = BytesIO()
                img.save(buf, "BMP")
                files[f"page_{i+1}.bmp"] = buf.getvalue()
            doc.close()
            log_conversion(request, "pdf_to_bmp")
            return zip_response(files, "pages_bmp.zip")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → BMP",
        "description": "Convert PDF pages to BMP images.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "🖼️",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → TIFF
# ---------------------------------------------------------------------------

def pdf_to_tiff(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            import fitz
            from PIL import Image
            uploaded = request.FILES["pdf_file"]
            doc = fitz.open(stream=uploaded.read(), filetype="pdf")
            files = {}
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                buf = BytesIO()
                img.save(buf, "TIFF")
                files[f"page_{i+1}.tiff"] = buf.getvalue()
            doc.close()
            log_conversion(request, "pdf_to_tiff")
            return zip_response(files, "pages_tiff.zip")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → TIFF",
        "description": "Convert PDF pages to TIFF images.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "🖼️",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → PPM
# ---------------------------------------------------------------------------

def pdf_to_ppm(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            import fitz
            uploaded = request.FILES["pdf_file"]
            doc = fitz.open(stream=uploaded.read(), filetype="pdf")
            files = {}
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                files[f"page_{i+1}.ppm"] = pix.tobytes("ppm")
            doc.close()
            log_conversion(request, "pdf_to_ppm")
            return zip_response(files, "pages_ppm.zip")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → PPM",
        "description": "Convert PDF pages to PPM image format.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "🖼️",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → SVG
# ---------------------------------------------------------------------------

def pdf_to_svg(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            import fitz
            uploaded = request.FILES["pdf_file"]
            doc = fitz.open(stream=uploaded.read(), filetype="pdf")
            files = {}
            for i, page in enumerate(doc):
                svg_text = page.get_svg_image(matrix=fitz.Matrix(1, 1))
                files[f"page_{i+1}.svg"] = svg_text.encode("utf-8")
            doc.close()
            log_conversion(request, "pdf_to_svg")
            return zip_response(files, "pages_svg.zip")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → SVG",
        "description": "Convert PDF pages to SVG vector graphics.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "📐",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → PPTX  (each page as image slide)
# ---------------------------------------------------------------------------

def pdf_to_pptx(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            import fitz
            from pptx import Presentation
            from pptx.util import Emu
            uploaded = request.FILES["pdf_file"]
            doc = fitz.open(stream=uploaded.read(), filetype="pdf")
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            for page in doc:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_buf = BytesIO(pix.tobytes("png"))
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    img_buf, 0, 0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )
            doc.close()
            buf = BytesIO()
            prs.save(buf)
            buf.seek(0)
            log_conversion(request, "pdf_to_pptx")
            return file_response(
                buf.read(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "converted.pptx",
            )
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → PPTX",
        "description": "Turn each PDF page into a PowerPoint slide.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "📊",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → XLSX
# ---------------------------------------------------------------------------

def pdf_to_xlsx(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            import pandas as pd
            from pypdf import PdfReader
            uploaded = request.FILES["pdf_file"]
            reader = PdfReader(uploaded)
            rows = [
                {"Page": i + 1, "Content": (page.extract_text() or "").strip()}
                for i, page in enumerate(reader.pages)
            ]
            df = pd.DataFrame(rows)
            buf = BytesIO()
            with pd.ExcelWriter(buf, engine="xlsxwriter") as writer:
                df.to_excel(writer, index=False)
            buf.seek(0)
            log_conversion(request, "pdf_to_xlsx")
            return file_response(
                buf.read(),
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "extracted.xlsx",
            )
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → XLSX",
        "description": "Extract PDF text into an Excel spreadsheet, one row per page.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "📊",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → CSV
# ---------------------------------------------------------------------------

def pdf_to_csv(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            from pypdf import PdfReader
            uploaded = request.FILES["pdf_file"]
            reader = PdfReader(uploaded)
            buf = io.StringIO()
            writer = csv.writer(buf)
            writer.writerow(["Page", "Content"])
            for i, page in enumerate(reader.pages):
                writer.writerow([i + 1, (page.extract_text() or "").strip()])
            log_conversion(request, "pdf_to_csv")
            return file_response(
                buf.getvalue().encode("utf-8"), "text/csv", "extracted.csv"
            )
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → CSV",
        "description": "Extract PDF text content into a CSV file.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "📋",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → ODT
# ---------------------------------------------------------------------------

def pdf_to_odt(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            from pypdf import PdfReader
            from odf.opendocument import OpenDocumentText
            from odf.text import P
            uploaded = request.FILES["pdf_file"]
            reader = PdfReader(uploaded)
            text = "\n\n".join(p.extract_text() or "" for p in reader.pages)
            doc = OpenDocumentText()
            for line in text.split("\n"):
                doc.text.addElement(P(text=line))
            with TempDir() as tmp:
                out = tmp / "output.odt"
                doc.save(str(out))
                log_conversion(request, "pdf_to_odt")
                return file_response(
                    out.read_bytes(),
                    "application/vnd.oasis.opendocument.text",
                    "converted.odt",
                )
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → ODT",
        "description": "Convert a PDF to OpenDocument Text format.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "📝",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# PDF → ZIP  (wrap PDF in a zip archive)
# ---------------------------------------------------------------------------

def pdf_to_zip(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            uploaded = request.FILES["pdf_file"]
            name = uploaded.name or "document.pdf"
            log_conversion(request, "pdf_to_zip")
            return zip_response({name: uploaded.read()}, "document.zip")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "PDF → ZIP",
        "description": "Package your PDF file inside a ZIP archive.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "icon": "🗜️",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# Split PDF
# ---------------------------------------------------------------------------

def pdf_split(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            from pypdf import PdfReader, PdfWriter
            uploaded = request.FILES["pdf_file"]
            reader = PdfReader(uploaded)
            total_pages = len(reader.pages)

            # Get how many pages per split chunk (default 1 = every page)
            try:
                pages_per_split = int(request.POST.get("pages_per_split", 1))
            except (ValueError, TypeError):
                pages_per_split = 1
            pages_per_split = max(1, min(pages_per_split, total_pages))

            files = {}
            chunk_num = 1
            for start in range(0, total_pages, pages_per_split):
                end = min(start + pages_per_split, total_pages)
                writer = PdfWriter()
                for i in range(start, end):
                    writer.add_page(reader.pages[i])
                buf = BytesIO()
                writer.write(buf)
                if pages_per_split == 1:
                    files[f"page_{chunk_num}.pdf"] = buf.getvalue()
                else:
                    files[f"pages_{start+1}-{end}.pdf"] = buf.getvalue()
                chunk_num += 1
            log_conversion(request, "pdf_split")
            return zip_response(files, "split_pages.zip")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "Split PDF",
        "description": "Split a PDF into smaller PDF files by choosing how many pages per file.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "extra_fields": [{
            "name": "pages_per_split",
            "label": "Pages per split file",
            "type": "select",
            "options": [
                {"value": "1", "text": "1 page per file (every page separate)"},
                {"value": "2", "text": "2 pages per file"},
                {"value": "3", "text": "3 pages per file"},
                {"value": "4", "text": "4 pages per file"},
                {"value": "5", "text": "5 pages per file"},
                {"value": "6", "text": "6 pages per file"},
                {"value": "7", "text": "7 pages per file"},
                {"value": "8", "text": "8 pages per file"},
                {"value": "9", "text": "9 pages per file"},
                {"value": "10", "text": "10 pages per file"},
            ],
        }],
        "icon": "✂️",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# Merge PDFs
# ---------------------------------------------------------------------------

def pdf_merge(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            from pypdf import PdfWriter
            files = request.FILES.getlist("pdf_files")
            if len(files) < 2:
                raise ValueError("Please upload at least 2 PDF files to merge.")
            writer = PdfWriter()
            for f in files:
                writer.append(f)
            buf = BytesIO()
            writer.write(buf)
            buf.seek(0)
            log_conversion(request, "pdf_merge")
            return file_response(buf.read(), "application/pdf", "merged.pdf")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "Merge PDFs",
        "description": "Combine multiple PDF files into one document.",
        "accept": ".pdf",
        "field_name": "pdf_files",
        "multiple": True,
        "icon": "🔗",
        "category_color": "red",
    })


# ---------------------------------------------------------------------------
# Encrypt PDF
# ---------------------------------------------------------------------------

def pdf_encrypt(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            from pypdf import PdfReader, PdfWriter
            uploaded = request.FILES["pdf_file"]
            password = request.POST.get("password", "").strip()
            if not password:
                raise ValueError("Please provide a password.")
            reader = PdfReader(uploaded)
            writer = PdfWriter()
            writer.append_pages_from_reader(reader)
            writer.encrypt(user_password=password, owner_password=password)
            buf = BytesIO()
            writer.write(buf)
            buf.seek(0)
            log_conversion(request, "pdf_encrypt")
            return file_response(buf.read(), "application/pdf", "encrypted.pdf")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "Encrypt PDF",
        "description": "Password-protect your PDF with AES encryption.",
        "accept": ".pdf",
        "field_name": "pdf_file",
        "extra_fields": [{"name": "password", "label": "Password", "type": "password"}],
        "icon": "🔒",
        "category_color": "red",
    })
