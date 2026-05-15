"""DOCX conversion tools – fully pure Python, no LibreOffice."""
from io import BytesIO

from django.shortcuts import render
from tools.utils import (
    TempDir, file_response, zip_response,
    docx_to_pdf_bytes,
)
from tools.conversion_limit import check_conversion_limit, log_conversion


def _check_limit(request):
    allowed, remaining = check_conversion_limit(request)
    if not allowed:
        return render(request, "tools/error.html", {
            "error": "You have used all 3 free conversions. Please sign up for unlimited access."
        })
    return None


def docx_to_pdf(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            uploaded = request.FILES["docx_file"]
            data = docx_to_pdf_bytes(uploaded)
            log_conversion(request, "docx_to_pdf")
            return file_response(data, "application/pdf", "converted.pdf")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "DOCX → PDF",
        "description": "Convert a Word document to PDF.",
        "accept": ".docx",
        "field_name": "docx_file",
        "icon": "📄",
        "category_color": "blue",
    })


def docx_to_txt(request):
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            from docx import Document
            uploaded = request.FILES["docx_file"]
            doc = Document(uploaded)
            text = "\n".join(p.text for p in doc.paragraphs)
            log_conversion(request, "docx_to_txt")
            return file_response(text.encode("utf-8"), "text/plain", "extracted.txt")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "DOCX → TXT",
        "description": "Extract all text from a Word document.",
        "accept": ".docx",
        "field_name": "docx_file",
        "icon": "📃",
        "category_color": "blue",
    })


def docx_to_odt(request):
    """DOCX → ODT using python-docx to read + odfpy to write."""
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            from docx import Document
            from odf.opendocument import OpenDocumentText
            from odf.text import P, H
            from odf.style import Style, TextProperties, ParagraphProperties

            uploaded = request.FILES["docx_file"]
            doc = Document(uploaded)

            odt = OpenDocumentText()
            for para in doc.paragraphs:
                style_name = para.style.name if para.style else ""
                if style_name.startswith("Heading"):
                    try:
                        lvl = int(style_name.split()[-1])
                    except ValueError:
                        lvl = 1
                    elem = H(outlinelevel=lvl, text=para.text)
                else:
                    elem = P(text=para.text)
                odt.text.addElement(elem)

            with TempDir() as tmp:
                out = tmp / "converted.odt"
                odt.save(str(out))
                log_conversion(request, "docx_to_odt")
                return file_response(
                    out.read_bytes(),
                    "application/vnd.oasis.opendocument.text",
                    "converted.odt",
                )
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "DOCX → ODT",
        "description": "Convert a Word document to OpenDocument Text format.",
        "accept": ".docx",
        "field_name": "docx_file",
        "icon": "📝",
        "category_color": "blue",
    })


def docx_to_pptx(request):
    """DOCX → PPTX: each paragraph/heading becomes a slide bullet."""
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            from docx import Document
            from pptx import Presentation
            from pptx.util import Inches, Pt

            uploaded = request.FILES["docx_file"]
            doc = Document(uploaded)
            prs = Presentation()
            blank = prs.slide_layouts[1]  # title + content

            current_title = None
            current_bullets = []

            def _flush():
                if current_title is None and not current_bullets:
                    return
                slide = prs.slides.add_slide(blank)
                slide.shapes.title.text = current_title or ""
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, b in enumerate(current_bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = b

            for para in doc.paragraphs:
                style_name = para.style.name if para.style else ""
                text = para.text.strip()
                if not text:
                    continue
                if style_name.startswith("Heading"):
                    _flush()
                    current_title = text
                    current_bullets = []
                else:
                    current_bullets.append(text)
                    # Max ~8 bullets per slide then auto-break
                    if len(current_bullets) >= 8:
                        _flush()
                        current_title = (current_title or "") + " (cont.)"
                        current_bullets = []

            _flush()
            if len(prs.slides) == 0:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = "Converted from DOCX"

            buf = BytesIO()
            prs.save(buf)
            buf.seek(0)
            log_conversion(request, "docx_to_pptx")
            return file_response(
                buf.read(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "converted.pptx",
            )
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "DOCX → PPTX",
        "description": "Convert Word headings and paragraphs into PowerPoint slides.",
        "accept": ".docx",
        "field_name": "docx_file",
        "icon": "📊",
        "category_color": "blue",
    })


def docx_to_png(request):
    """DOCX → PNG: convert via in-memory PDF then render with pymupdf."""
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            import fitz
            uploaded = request.FILES["docx_file"]
            pdf_bytes = docx_to_pdf_bytes(uploaded)
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            files = {}
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                files[f"page_{i+1}.png"] = pix.tobytes("png")
            doc.close()
            log_conversion(request, "docx_to_png")
            return zip_response(files, "pages_png.zip")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "DOCX → PNG",
        "description": "Convert each Word document page into a PNG image.",
        "accept": ".docx",
        "field_name": "docx_file",
        "icon": "🖼️",
        "category_color": "blue",
    })


def docx_to_jpg(request):
    """DOCX → JPG: convert via in-memory PDF then render with pymupdf."""
    if request.method == "POST":
        limit_resp = _check_limit(request)
        if limit_resp:
            return limit_resp
        try:
            import fitz
            uploaded = request.FILES["docx_file"]
            pdf_bytes = docx_to_pdf_bytes(uploaded)
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            files = {}
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                files[f"page_{i+1}.jpg"] = pix.tobytes("jpeg")
            doc.close()
            log_conversion(request, "docx_to_jpg")
            return zip_response(files, "pages_jpg.zip")
        except Exception as e:
            return render(request, "tools/error.html", {"error": str(e)})
    return render(request, "tools/tool_page.html", {
        "title": "DOCX → JPG",
        "description": "Convert each Word document page into a JPEG image.",
        "accept": ".docx",
        "field_name": "docx_file",
        "icon": "🖼️",
        "category_color": "blue",
    })
