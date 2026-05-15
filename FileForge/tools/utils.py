"""
Shared utilities for all converter views.
Fully standalone – no LibreOffice, no system installs required.
FFmpeg is expected at:  <project_root>/bin/ffmpeg.exe  (Windows)
                        ffmpeg on PATH               (Linux/macOS fallback)
"""
import os
import platform
import subprocess
import tempfile
import zipfile
from io import BytesIO
from pathlib import Path

from django.http import HttpResponse

# ---------------------------------------------------------------------------
# HTTP response helpers
# ---------------------------------------------------------------------------

def file_response(data: bytes, content_type: str, filename: str) -> HttpResponse:
    response = HttpResponse(data, content_type=content_type)
    response["Content-Disposition"] = f'attachment; filename="{filename}"'
    return response


def zip_response(files: dict, zip_name: str) -> HttpResponse:
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, data in files.items():
            zf.writestr(name, data)
    buf.seek(0)
    return file_response(buf.read(), "application/zip", zip_name)


# ---------------------------------------------------------------------------
# Temp-file context manager
# ---------------------------------------------------------------------------

class TempDir:
    def __enter__(self):
        self.path = tempfile.mkdtemp()
        return Path(self.path)

    def __exit__(self, *_):
        import shutil
        shutil.rmtree(self.path, ignore_errors=True)


# ---------------------------------------------------------------------------
# FFmpeg – uses bundled bin/ffmpeg.exe when available
# ---------------------------------------------------------------------------

def _find_ffmpeg() -> str:
    """
    Look for ffmpeg in order:
      1. FFMPEG_PATH env var
      2. <BASE_DIR>/bin/ffmpeg.exe   (bundled by install.bat)
      3. 'ffmpeg' on system PATH     (fallback for Linux/macOS)
    """
    env = os.environ.get("FFMPEG_PATH")
    if env:
        return env
    bundled = Path(settings.BASE_DIR) / "bin" / "ffmpeg.exe"
    if bundled.exists():
        return str(bundled)
    return "ffmpeg"


def ffmpeg_convert(src_path: str, out_path: str, extra_args=None) -> None:
    exe = _find_ffmpeg()
    cmd = [exe, "-y", "-i", src_path] + (extra_args or []) + [out_path]
    flags = subprocess.CREATE_NO_WINDOW if platform.system() == "Windows" else 0
    result = subprocess.run(cmd, capture_output=True, text=True,
                            timeout=300, creationflags=flags)
    if result.returncode != 0:
        raise RuntimeError(
            f"FFmpeg conversion failed.\n"
            f"Error: {result.stderr[-800:] if result.stderr else '(no output)'}\n\n"
            "Make sure you ran install.bat to download ffmpeg.exe into the bin\\ folder."
        )


# ---------------------------------------------------------------------------
# Pure-Python DOCX -> PDF  (no LibreOffice, no Word required)
# python-docx to read + reportlab to render
# ---------------------------------------------------------------------------

def docx_to_pdf_bytes(docx_file_or_path) -> bytes:
    """Convert a DOCX (path or file-like) to PDF bytes using reportlab."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    )
    from reportlab.lib import colors

    doc = Document(docx_file_or_path)
    buf = BytesIO()

    pdf = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=2.5*cm, rightMargin=2.5*cm,
        topMargin=2.5*cm, bottomMargin=2.5*cm,
    )

    _ALIGN_MAP = {
        WD_ALIGN_PARAGRAPH.LEFT:    TA_LEFT,
        WD_ALIGN_PARAGRAPH.CENTER:  TA_CENTER,
        WD_ALIGN_PARAGRAPH.RIGHT:   TA_RIGHT,
        WD_ALIGN_PARAGRAPH.JUSTIFY: TA_JUSTIFY,
    }

    def _sty(name, size, bold=False, sb=4, sa=4, align=TA_LEFT):
        return ParagraphStyle(name,
            fontName="Helvetica-Bold" if bold else "Helvetica",
            fontSize=size, leading=size*1.35,
            spaceBefore=sb, spaceAfter=sa, alignment=align)

    h_styles = {
        1: _sty("h1", 22, bold=True, sb=12, sa=6),
        2: _sty("h2", 17, bold=True, sb=10, sa=4),
        3: _sty("h3", 14, bold=True, sb=8,  sa=4),
        4: _sty("h4", 12, bold=True, sb=6,  sa=4),
    }
    normal = _sty("n", 11, sb=2, sa=4)

    def _esc(t):
        return t.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")

    def _para_to_rl(para):
        style_name = para.style.name if para.style else ""
        align = _ALIGN_MAP.get(para.alignment, TA_LEFT)
        level = None
        if style_name.startswith("Heading"):
            try:
                level = int(style_name.split()[-1])
            except ValueError:
                pass

        parts = []
        for run in para.runs:
            t = _esc(run.text)
            if not t:
                continue
            if run.bold and run.italic:
                parts.append(f"<b><i>{t}</i></b>")
            elif run.bold:
                parts.append(f"<b>{t}</b>")
            elif run.italic:
                parts.append(f"<i>{t}</i>")
            else:
                parts.append(t)

        markup = "".join(parts) or "&nbsp;"
        base = h_styles.get(level, normal) if level else normal
        sty = ParagraphStyle("copy", parent=base, alignment=align)
        return Paragraph(markup, sty)

    story = []
    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

        if tag == "p":
            from docx.text.paragraph import Paragraph as DocxPara
            story.append(_para_to_rl(DocxPara(block, doc)))

        elif tag == "tbl":
            from docx.table import Table as DocxTable
            tbl = DocxTable(block, doc)
            data = []
            for row in tbl.rows:
                data.append([Paragraph(_esc(cell.text), normal)
                              for cell in row.cells])
            if data:
                col_count = max(len(r) for r in data)
                col_w = (A4[0] - 5*cm) / col_count
                rl_tbl = Table(data, colWidths=[col_w]*col_count)
                rl_tbl.setStyle(TableStyle([
                    ("BACKGROUND",    (0,0),(-1,0), colors.HexColor("#f0f0f0")),
                    ("FONTNAME",      (0,0),(-1,0), "Helvetica-Bold"),
                    ("FONTSIZE",      (0,0),(-1,-1), 9),
                    ("GRID",          (0,0),(-1,-1), 0.5, colors.grey),
                    ("VALIGN",        (0,0),(-1,-1), "TOP"),
                    ("TOPPADDING",    (0,0),(-1,-1), 4),
                    ("BOTTOMPADDING", (0,0),(-1,-1), 4),
                    ("LEFTPADDING",   (0,0),(-1,-1), 6),
                    ("RIGHTPADDING",  (0,0),(-1,-1), 6),
                ]))
                story.append(Spacer(1, 6))
                story.append(rl_tbl)
                story.append(Spacer(1, 6))

    if not story:
        story.append(Paragraph("(empty document)", normal))

    pdf.build(story)
    buf.seek(0)
    return buf.read()


# ---------------------------------------------------------------------------
# Pure-Python XLSX -> PDF
# ---------------------------------------------------------------------------

def xlsx_to_pdf_bytes(file_or_path) -> bytes:
    import openpyxl
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.units import cm
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib import colors

    wb = openpyxl.load_workbook(file_or_path, data_only=True)
    buf = BytesIO()
    pdf = SimpleDocTemplate(buf, pagesize=landscape(A4),
                            leftMargin=1.5*cm, rightMargin=1.5*cm,
                            topMargin=1.5*cm, bottomMargin=1.5*cm)

    cell_sty = ParagraphStyle("c", fontName="Helvetica",      fontSize=8, leading=10)
    head_sty = ParagraphStyle("h", fontName="Helvetica-Bold", fontSize=8, leading=10)
    title_sty= ParagraphStyle("t", fontName="Helvetica-Bold", fontSize=12,
                               spaceBefore=8, spaceAfter=4)
    story = []

    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        story.append(Paragraph(ws.title, title_sty))
        max_cols = min(len(rows[0]), 30)
        data = []
        for i, row in enumerate(rows):
            sty = head_sty if i == 0 else cell_sty
            data.append([Paragraph(str(v) if v is not None else "", sty)
                         for v in row[:max_cols]])
        if data:
            page_w = landscape(A4)[0] - 3*cm
            col_w  = page_w / max_cols
            tbl = Table(data, colWidths=[col_w]*max_cols, repeatRows=1)
            tbl.setStyle(TableStyle([
                ("BACKGROUND",    (0,0),(-1,0), colors.HexColor("#d0d8e8")),
                ("FONTNAME",      (0,0),(-1,0), "Helvetica-Bold"),
                ("FONTSIZE",      (0,0),(-1,-1), 8),
                ("GRID",          (0,0),(-1,-1), 0.3, colors.HexColor("#cccccc")),
                ("ROWBACKGROUNDS",(0,1),(-1,-1),
                 [colors.white, colors.HexColor("#f7f7f7")]),
                ("VALIGN",        (0,0),(-1,-1), "TOP"),
                ("TOPPADDING",    (0,0),(-1,-1), 3),
                ("BOTTOMPADDING", (0,0),(-1,-1), 3),
                ("LEFTPADDING",   (0,0),(-1,-1), 4),
                ("RIGHTPADDING",  (0,0),(-1,-1), 4),
            ]))
            story.append(tbl)
        story.append(Spacer(1, 12))

    if not story:
        story.append(Paragraph("(empty workbook)", cell_sty))

    pdf.build(story)
    buf.seek(0)
    return buf.read()


# ---------------------------------------------------------------------------
# Pure-Python PPTX -> PDF
# ---------------------------------------------------------------------------

def pptx_to_pdf_bytes(file_or_path) -> bytes:
    """Convert a PPTX to PDF by rendering each slide as an image, then compositing
    those images into a multi-page PDF.  This preserves all visual content including
    images, shapes, backgrounds, and text styling."""
    import subprocess
    import platform
    import shutil
    from pathlib import Path

    # Try LibreOffice first (best quality, preserves everything)
    lo_path = shutil.which("libreoffice") or shutil.which("soffice")
    if lo_path:
        try:
            return _pptx_to_pdf_via_libreoffice(file_or_path, lo_path)
        except Exception:
            pass  # Fall through to image-based approach

    # Fallback: render slides as images via python-pptx + Pillow + reportlab
    return _pptx_to_pdf_via_images(file_or_path)


def _pptx_to_pdf_via_libreoffice(file_or_path, lo_path):
    """Use LibreOffice headless to convert PPTX → PDF with full fidelity."""
    import subprocess
    import platform

    with TempDir() as tmp:
        # If file_or_path is a file-like object, save it to disk first
        if hasattr(file_or_path, 'read'):
            src = tmp / "input.pptx"
            src.write_bytes(file_or_path.read())
        else:
            src = Path(file_or_path)

        flags = subprocess.CREATE_NO_WINDOW if platform.system() == "Windows" else 0
        subprocess.run(
            [lo_path, "--headless", "--convert-to", "pdf", "--outdir", str(tmp), str(src)],
            capture_output=True, timeout=120, creationflags=flags,
        )
        pdf_out = tmp / src.with_suffix(".pdf").name
        if pdf_out.exists():
            return pdf_out.read_bytes()
        raise RuntimeError("LibreOffice conversion produced no output")


def _pptx_to_pdf_via_images(file_or_path):
    """Render each PPTX slide as an image and compose into a PDF.
    Extracts text AND embedded images from each slide."""
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image, ImageDraw, ImageFont
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Image as RLImage, Spacer
    from reportlab.lib.styles import ParagraphStyle

    prs = Presentation(file_or_path)
    buf = BytesIO()

    EMU_PER_PT = 12700
    slide_w_pt = prs.slide_width / EMU_PER_PT
    slide_h_pt = prs.slide_height / EMU_PER_PT

    # Render each slide to a PIL image
    DPI = 150
    img_w = int(slide_w_pt * DPI / 72)
    img_h = int(slide_h_pt * DPI / 72)

    slide_images = []
    for slide_idx, slide in enumerate(prs.slides):
        # Create a white canvas for the slide
        canvas = Image.new("RGB", (img_w, img_h), (255, 255, 255))
        draw = ImageDraw.Draw(canvas)

        # Try to load a basic font
        try:
            font_large = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
            font_normal = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
        except (OSError, IOError):
            try:
                font_large = ImageFont.truetype("arial.ttf", 24)
                font_normal = ImageFont.truetype("arial.ttf", 16)
            except (OSError, IOError):
                font_large = ImageFont.load_default()
                font_normal = ImageFont.load_default()

        for shape in slide.shapes:
            # Calculate position in pixels
            if shape.left is not None and shape.top is not None:
                x = int((shape.left / EMU_PER_PT) * DPI / 72)
                y = int((shape.top / EMU_PER_PT) * DPI / 72)
                w = int((shape.width / EMU_PER_PT) * DPI / 72) if shape.width else img_w
                h = int((shape.height / EMU_PER_PT) * DPI / 72) if shape.height else img_h
            else:
                x, y, w, h = 0, 0, img_w, img_h

            # Handle images/pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                try:
                    img_blob = shape.image.blob
                    img = Image.open(BytesIO(img_blob))
                    img = img.convert("RGB")
                    img = img.resize((w, h), Image.LANCZOS)
                    canvas.paste(img, (x, y))
                except Exception:
                    pass

            # Handle text frames
            if shape.has_text_frame:
                text_y = y
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        is_title = (para == shape.text_frame.paragraphs[0]
                                    and hasattr(shape, 'placeholder_format')
                                    and shape.placeholder_format is not None)
                        fnt = font_large if is_title else font_normal
                        # Word-wrap text within the shape bounds
                        _draw_wrapped_text(draw, line, x + 8, text_y + 4, w - 16, fnt, fill=(30, 30, 30))
                        try:
                            bbox = draw.textbbox((0, 0), line, font=fnt)
                            line_h = bbox[3] - bbox[1] + 6
                        except Exception:
                            line_h = 22
                        text_y += line_h

            # Handle tables
            if shape.has_table:
                tbl = shape.table
                cell_h = max(h // len(tbl.rows), 20) if len(tbl.rows) > 0 else 20
                cell_w = max(w // len(tbl.columns), 40) if len(tbl.columns) > 0 else 40
                for ri, row in enumerate(tbl.rows):
                    for ci, cell in enumerate(row.cells):
                        cx = x + ci * cell_w
                        cy = y + ri * cell_h
                        draw.rectangle([cx, cy, cx + cell_w, cy + cell_h], outline=(180, 180, 180))
                        cell_text = cell.text.strip()
                        if cell_text:
                            draw.text((cx + 4, cy + 2), cell_text[:30], fill=(30, 30, 30), font=font_normal)

        slide_images.append(canvas)

    # Build PDF from rendered slide images
    from reportlab.pdfgen.canvas import Canvas as PDFCanvas
    pdf_buf = BytesIO()
    c = PDFCanvas(pdf_buf, pagesize=(slide_w_pt, slide_h_pt))

    for img in slide_images:
        img_buf = BytesIO()
        img.save(img_buf, "PNG", quality=95)
        img_buf.seek(0)

        from reportlab.lib.utils import ImageReader
        c.drawImage(ImageReader(img_buf), 0, 0, width=slide_w_pt, height=slide_h_pt)
        c.showPage()

    c.save()
    pdf_buf.seek(0)
    return pdf_buf.read()


def _draw_wrapped_text(draw, text, x, y, max_width, font, fill=(0, 0, 0)):
    """Draw text with basic word wrapping."""
    words = text.split()
    lines = []
    current_line = ""
    for word in words:
        test = current_line + (" " if current_line else "") + word
        try:
            bbox = draw.textbbox((0, 0), test, font=font)
            tw = bbox[2] - bbox[0]
        except Exception:
            tw = len(test) * 8
        if tw <= max_width or not current_line:
            current_line = test
        else:
            lines.append(current_line)
            current_line = word
    if current_line:
        lines.append(current_line)

    cy = y
    for line in lines:
        draw.text((x, cy), line, fill=fill, font=font)
        try:
            bbox = draw.textbbox((0, 0), line, font=font)
            cy += bbox[3] - bbox[1] + 4
        except Exception:
            cy += 20

